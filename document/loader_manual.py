import os, pandas as pd  # type: ignore
from typing import List

from langchain_google_genai import ChatGoogleGenerativeAI
from docling.document_converter import DocumentConverter
from langchain_core.documents import Document
from trafilatura import extract, fetch_url
from docx import Document as DocxDocument
from pptx import Presentation


class LoaderManual:
    @staticmethod
    def load_pdf(file_path: str) -> List[Document]:
        converter = DocumentConverter()
        result = converter.convert(file_path)
        return [
            Document(
                page_content=result.document.export_to_markdown(),
                metadata={"source": file_path, "type": "pdf"},
            )
        ]

    @staticmethod
    def load_html(path_or_url: str) -> List[Document]:
        downloaded = (
            fetch_url(path_or_url)
            if path_or_url.startswith("http")
            else open(path_or_url, "r").read()
        )
        content = extract(downloaded)
        return [Document(page_content=content or "", metadata={"source": path_or_url})]

    @staticmethod
    def load_docx(file_path: str) -> List[Document]:
        doc = DocxDocument(file_path)
        full_text = [para.text for para in doc.paragraphs]
        return [
            Document(page_content="\n".join(full_text), metadata={"source": file_path})
        ]

    @staticmethod
    def load_pptx(file_path: str) -> List[Document]:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return [
            Document(page_content="\n".join(text_runs), metadata={"source": file_path})
        ]

    @staticmethod
    def load_csv(file_path: str) -> List[Document]:
        df = pd.read_csv(file_path)
        return [Document(page_content=df.to_string(), metadata={"source": file_path})]

    @staticmethod
    async def load_image(file_path: str) -> List[Document]:
        llm = ChatGoogleGenerativeAI(model="gemini-1.5-flash")
        response = llm.invoke(
            [
                {
                    "type": "text",
                    "text": "Transcribe all text from this image accurately.",
                },
                {"type": "image_url", "image_url": file_path},
            ]
        )
        return [Document(page_content=response.content, metadata={"source": file_path})]

    @staticmethod
    def load_text(file_path: str) -> List[Document]:
        with open(file_path, "r") as f:
            return [Document(page_content=f.read(), metadata={"source": file_path})]

    @classmethod
    def process_file(cls, file_path: str) -> List[Document]:
        ext = os.path.splitext(file_path)[-1].lower()
        print(f"🔍 Extracting {ext}: {file_path}")

        match ext:
            case ".pdf":
                return cls.load_pdf(file_path)
            case ".docx":
                return cls.load_docx(file_path)
            case ".pptx":
                return cls.load_pptx(file_path)
            case ".csv":
                return cls.load_csv(file_path)
            case ".html":
                return cls.load_html(file_path)
            case ".jpg", ".jpeg", ".png":
                return cls.load_image(file_path)
            case ".txt", ".md":
                return cls.load_text(file_path)
            case _:
                raise ValueError(f"Unsupported file type: {ext}")
