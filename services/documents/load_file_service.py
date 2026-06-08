from models import (
    FileNotFoundError,
    FileTypeError,
    DocumentLoadError,
    get_checked_env_var,
)
from langchain_google_genai import ChatGoogleGenerativeAI
from docling.document_converter import DocumentConverter
from langchain_unstructured import UnstructuredLoader
from langchain_core.messages import HumanMessage
from langchain_core.documents import Document
from trafilatura import extract, fetch_url
from docx import Document as DocxDocument
from pptx import Presentation
import logging, base64, json, os
from typing import List
import pandas  # pyright: ignore

logger = logging.getLogger(__name__)


def load_pdf(file: str) -> List[Document]:
    converter = DocumentConverter()
    result = converter.convert(file)

    return [
        Document(
            page_content=result.document.export_to_markdown(),
            metadata={"source": file, "type": "pdf"},
        )
    ]


def load_html(path_or_url: str) -> List[Document]:
    downloaded = (
        fetch_url(path_or_url)
        if path_or_url.startswith("http")
        else open(path_or_url, "r", encoding="utf-8").read()
    )

    content = extract(downloaded)

    return [
        Document(
            page_content=content or "",
            metadata={"source": path_or_url},
        )
    ]


def load_docx(file: str) -> List[Document]:
    doc = DocxDocument(file)

    return [
        Document(
            page_content="\n".join([p.text for p in doc.paragraphs]),
            metadata={"source": file},
        )
    ]


def load_pptx(file: str) -> List[Document]:
    prs = Presentation(file)
    documents = []

    def get_shape_text(shape) -> List[str]:
        parts = []

        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    parts.append(text)

        elif shape.has_table:
            for row in shape.table.rows:
                row_text = [
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                ]

                if row_text:
                    parts.append(" | ".join(row_text))

        elif hasattr(shape, "shape_type") and shape.shape_type == 6:
            for sub_shape in shape.shapes:
                parts.extend(get_shape_text(sub_shape))

        return parts

    for i, slide in enumerate(prs.slides):
        slide_content = []

        for shape in slide.shapes:
            slide_content.extend(get_shape_text(shape))

        if slide_content:
            documents.append(
                Document(
                    page_content="\n".join(slide_content),
                    metadata={
                        "source": file,
                        "slide_number": i + 1,
                    },
                )
            )

    return documents


def load_csv(file: str) -> List[Document]:
    df = pandas.read_csv(file)

    return [
        Document(
            page_content=df.to_string(),
            metadata={"source": file},
        )
    ]


def load_image(file: str) -> List[Document]:
    api_key = get_checked_env_var("GOOGLE_API_KEY", logger)
    llm = ChatGoogleGenerativeAI(
        model="gemini-2.5-flash",
        google_api_key=api_key,
    )

    with open(file, "rb") as f:
        image_b64 = base64.b64encode(f.read()).decode("utf-8")

    message = HumanMessage(
        content=[
            {
                "type": "text",
                "text": "Transcribe all text from this image accurately.",
            },
            {
                "type": "image_url",
                "image_url": {"url": f"data:image/png;base64,{image_b64}"},
            },
        ]
    )

    response = llm.invoke([message])

    return [
        Document(
            page_content=str(response.content),
            metadata={
                "source": file,
                "type": "image",
            },
        )
    ]


def load_text(file: str) -> List[Document]:
    with open(file, "r", encoding="utf-8") as f:
        return [
            Document(
                page_content=f.read(),
                metadata={"source": file},
            )
        ]


def write_text(loader_type: str, file: str, content: str) -> None:
    output_dir = f"data/extracts/{loader_type}"
    os.makedirs(output_dir, exist_ok=True)

    with open(
        f"{output_dir}/{file}.txt",
        "w",
        encoding="utf-8",
    ) as f:
        json.dump(content, f, indent=4)


def check_file_exists(file: str) -> None:
    if not os.path.exists(file):
        msg = f"❌ File {file} not found"
        logger.error(msg)
        raise FileNotFoundError(msg)


def load_manual(file: str) -> List[Document]:
    check_file_exists(file)
    ext = os.path.splitext(file)[-1].lower()

    try:
        match ext:
            case ".pdf":
                return load_pdf(file)

            case ".docx":
                return load_docx(file)

            case ".pptx":
                return load_pptx(file)

            case ".csv":
                return load_csv(file)

            case ".html":
                return load_html(file)

            case ".jpg" | ".jpeg" | ".png":
                return load_image(file)

            case ".txt" | ".md":
                return load_text(file)

            case _:
                msg = f"❌ Unsupported file type {ext}"
                logger.error(msg)
                raise FileTypeError(msg)

    except Exception as e:
        msg = f"❌ Failed to load file {file} using manual loader: {str(e)}"
        logger.exception(msg)
        raise DocumentLoadError(msg) from e


def load_unstructured(file: str) -> List[Document]:
    check_file_exists(file)
    api_key = get_checked_env_var("UNSTRUCTURED_API_KEY", logger)

    try:
        loader = UnstructuredLoader(
            file_path=[file],
            mode="elements",
            partition_via_api=True,
            api_key=api_key,
        )

        return loader.load()

    except Exception as e:
        msg = f"❌ Failed to load file {file} using Unstructured loader: {str(e)}"
        logger.exception(msg)
        raise DocumentLoadError(msg) from e
