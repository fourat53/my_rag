import os, pandas as pd
from typing import List

from langchain_google_genai import ChatGoogleGenerativeAI
from docling.document_converter import DocumentConverter
from langchain_core.documents import Document
from trafilatura import extract, fetch_url
from docx import Document as DocxDocument
from pptx import Presentation


class LoadManual:
    @staticmethod
    def load_pdf(file_path: str) -> List[Document]:
        converter = DocumentConverter()
        result = converter.convert(file_path)
        return [
            Document(
                page_content=result.document.export_to_markdown(),
                metadata={"source": file_path, "type": "pdf"},
            )
        ]

    @staticmethod
    def load_html(path_or_url: str) -> List[Document]:
        downloaded = (
            fetch_url(path_or_url)
            if path_or_url.startswith("http")
            else open(path_or_url, "r").read()
        )
        content = extract(downloaded)
        return [Document(page_content=content or "", metadata={"source": path_or_url})]

    @staticmethod
    def load_docx(file_path: str) -> List[Document]:
        doc = DocxDocument(file_path)
        full_text = [para.text for para in doc.paragraphs]
        return [
            Document(page_content="\n".join(full_text), metadata={"source": file_path})
        ]

    @staticmethod
    def load_pptx(file_path: str) -> List[Document]:
        prs = Presentation(file_path)
        documents = []

        def get_shape_text(shape) -> List[str]:
            parts = []
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        parts.append(paragraph.text.strip())

            elif shape.has_table:
                for row in shape.table.rows:
                    row_text = [
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    ]
                    if row_text:
                        parts.append(" | ".join(row_text))

            elif hasattr(shape, "shape_type") and shape.shape_type == 6:
                for sub_shape in shape.shapes:
                    parts.extend(get_shape_text(sub_shape))

            return parts

        for i, slide in enumerate(prs.slides):
            slide_content = []
            for shape in slide.shapes:
                slide_content.extend(get_shape_text(shape))

            if slide_content:
                text = "\n".join(slide_content)
                documents.append(
                    Document(
                        page_content=text,
                        metadata={"source": file_path, "slide_number": i + 1},
                    )
                )

        return documents

    @staticmethod
    def load_csv(file_path: str) -> List[Document]:
        df = pd.read_csv(file_path)
        return [Document(page_content=df.to_string(), metadata={"source": file_path})]

    @staticmethod
    async def load_image(file_path: str) -> List[Document]:
        llm = ChatGoogleGenerativeAI(model="gemini-1.5-flash")

        message = [
            {
                "type": "text",
                "text": "Transcribe all text from this image accurately.",
            },
            {"type": "image_url", "image_url": file_path},
        ]

        response = llm.invoke(message)

        content = (
            str(response.content)
            if not isinstance(response.content, str)
            else response.content
        )

        return [Document(page_content=content, metadata={"source": file_path})]

    @staticmethod
    def load_text(file_path: str) -> List[Document]:
        with open(file_path, "r") as f:
            return [Document(page_content=f.read(), metadata={"source": file_path})]

    @classmethod
    def process_file(cls, file_path: str) -> List[Document]:
        relative_path = f"data/{file_path}"
        ext = os.path.splitext(file_path)[-1].lower()
        print(f"🔍 Extracting {ext}: {file_path}")
        match ext:
            case ".pdf":
                return cls.load_pdf(relative_path)
            case ".docx":
                return cls.load_docx(relative_path)
            case ".pptx":
                return cls.load_pptx(relative_path)
            case ".csv":
                return cls.load_csv(relative_path)
            case ".html":
                return cls.load_html(relative_path)
            case ".jpg", ".jpeg", ".png":
                return cls.load_image(relative_path)
            case ".txt", ".md":
                return cls.load_text(relative_path)
            case _:
                raise ValueError(f"Unsupported file type: {ext}")
